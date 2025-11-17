"""PowerPoint extraction helpers."""
from __future__ import annotations

from pathlib import Path
from typing import List

try:
    import pptx  # type: ignore
except ImportError:  # pragma: no cover
    pptx = None  # type: ignore


def extract_pptx(path: Path) -> str:
    """Extract slide text via python-pptx."""
    if pptx is None:
        raise RuntimeError("python-pptx not installed")
    prs = pptx.Presentation(str(path))
    texts: List[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts.append(f"[Slide {slide_idx}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.append(cell.text)
    return "\n".join(texts)


__all__ = ["extract_pptx"]
